"""Document loaders for GH Buddy."""

from __future__ import annotations

from pathlib import Path
from typing import List, Dict

import numpy as np
import requests
try:
    from bs4 import BeautifulSoup
except ModuleNotFoundError:
    BeautifulSoup = None

try:
    from docx import Document as DocxDocument
except ModuleNotFoundError:
    DocxDocument = None
try:
    from pptx import Presentation
except ModuleNotFoundError:
    Presentation = None

try:
    from PyPDF2 import PdfReader
except ModuleNotFoundError:
    PdfReader = None
try:
    import fitz  # PyMuPDF
except ModuleNotFoundError:
    fitz = None
try:
    import pypdfium2 as pdfium
except ModuleNotFoundError:
    pdfium = None
try:
    from rapidocr_onnxruntime import RapidOCR
except ModuleNotFoundError:
    RapidOCR = None


LoadedDoc = Dict[str, str]


def _is_text_quality_ok(text: str) -> bool:
    """Heuristic filter to skip very noisy OCR output."""
    cleaned = " ".join(text.split())
    if len(cleaned) < 80:
        return False
    letters = sum(ch.isalpha() for ch in cleaned)
    if letters / max(1, len(cleaned)) < 0.45:
        return False
    words = [w for w in cleaned.split(" ") if w]
    if len(words) < 12:
        return False
    avg_word_len = sum(len(w) for w in words) / len(words)
    if avg_word_len > 9.5:
        return False
    return True


def load_pdf(file_path: str) -> List[LoadedDoc]:
    """Load text from a PDF file."""
    docs: List[LoadedDoc] = []
    path = Path(file_path)

    # 1) Try PyMuPDF first (usually stronger extraction than PyPDF2).
    global fitz
    if fitz is None:
        try:
            import fitz as _fitz  # type: ignore

            fitz = _fitz
        except ModuleNotFoundError:
            fitz = None

    if fitz is not None:
        try:
            pdf = fitz.open(str(path))
            for page_idx, page in enumerate(pdf, start=1):
                page_text = (page.get_text("text") or "").strip()
                if page_text and _is_text_quality_ok(page_text):
                    docs.append({"content": page_text, "source": f"{path.name}, page {page_idx}"})
            pdf.close()
            if docs:
                return docs
        except Exception:
            docs = []

    # 2) Fallback to PyPDF2.
    global PdfReader
    if PdfReader is None:
        try:
            from PyPDF2 import PdfReader as _PdfReader

            PdfReader = _PdfReader
        except ModuleNotFoundError:
            return []

    try:
        reader = PdfReader(str(path))
        for page_idx, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text and _is_text_quality_ok(page_text):
                docs.append({"content": page_text, "source": f"{path.name}, page {page_idx}"})
    except Exception:
        docs = []

    if docs:
        return docs

    # 3) OCR fallback for scanned/image-only PDFs.
    if pdfium is not None and RapidOCR is not None:
        try:
            ocr = RapidOCR()
            pdf = pdfium.PdfDocument(str(path))
            for page_idx in range(len(pdf)):
                page = pdf[page_idx]
                bitmap = page.render(scale=2.0).to_pil()
                image_np = np.array(bitmap)
                result, _ = ocr(image_np)
                if result:
                    text = " ".join([line[1] for line in result if len(line) > 1 and line[1]])
                    text = text.strip()
                    if text and _is_text_quality_ok(text):
                        docs.append({"content": text, "source": f"{path.name}, page {page_idx + 1} (OCR)"})
            if docs:
                return docs
        except Exception:
            return []
    return docs


def load_docx(file_path: str) -> List[LoadedDoc]:
    """Load text from a DOCX file."""
    global DocxDocument
    if DocxDocument is None:
        try:
            from docx import Document as _DocxDocument

            DocxDocument = _DocxDocument
        except ModuleNotFoundError:
            return []
    docs: List[LoadedDoc] = []
    path = Path(file_path)
    document = DocxDocument(str(path))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    if paragraphs:
        docs.append({"content": "\n".join(paragraphs), "source": path.name})
    return docs


def load_pptx(file_path: str) -> List[LoadedDoc]:
    """Load text from a PPTX file slide by slide."""
    global Presentation
    if Presentation is None:
        try:
            from pptx import Presentation as _Presentation

            Presentation = _Presentation
        except ModuleNotFoundError:
            return []
    docs: List[LoadedDoc] = []
    path = Path(file_path)
    try:
        prs = Presentation(str(path))
    except Exception:
        return []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text_parts.append(shape.text.strip())
        slide_text = "\n".join([t for t in slide_text_parts if t]).strip()
        if slide_text and _is_text_quality_ok(slide_text):
            docs.append({"content": slide_text, "source": f"{path.name}, slide {idx}"})
    return docs


def load_txt(file_path: str) -> List[LoadedDoc]:
    """Load text from a plain text file."""
    path = Path(file_path)
    content = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not content:
        return []
    return [{"content": content, "source": path.name}]


def load_url(url: str, timeout: int = 15) -> List[LoadedDoc]:
    """Scrape text from a URL and return as a document."""
    global BeautifulSoup
    if BeautifulSoup is None:
        try:
            from bs4 import BeautifulSoup as _BeautifulSoup

            BeautifulSoup = _BeautifulSoup
        except ModuleNotFoundError:
            return []
    try:
        headers = {"User-Agent": "GH-Buddy/1.0 (+https://github.com/)"}
        response = requests.get(url, timeout=timeout, headers=headers)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = " ".join(soup.stripped_strings)
        if not text:
            return []
        return [{"content": text, "source": url}]
    except requests.RequestException:
        return []


def load_documents(file_paths: List[str] | None = None, urls: List[str] | None = None) -> List[LoadedDoc]:
    """Load mixed document types and return a unified list."""
    file_paths = file_paths or []
    urls = urls or []
    docs: List[LoadedDoc] = []

    for file_path in file_paths:
        suffix = Path(file_path).suffix.lower()
        if suffix == ".pdf":
            docs.extend(load_pdf(file_path))
        elif suffix == ".docx":
            docs.extend(load_docx(file_path))
        elif suffix == ".pptx":
            docs.extend(load_pptx(file_path))
        elif suffix == ".txt":
            docs.extend(load_txt(file_path))

    for url in urls:
        docs.extend(load_url(url))

    return docs
